#!/usr/bin/env python3
"""Build the Phase 02 (Quality Engineering) opening deck.

Reproduces the visual system of the Phase 00 deck
(marketing/decks/phase-00-opening-ai-dev.pptx) so the pair stays consistent:
16:9, Calibri Light / Calibri, blue #3D8DFF accent, gray #60646C, F6F6F6 cards.

Content is derived from docs/phases/02/{OVERVIEW,TALKING-POINTS,TEACHING}.md.

Usage: python3 scripts/build-phase-02-deck.py
Output: marketing/decks/phase-02-quality-engineering.pptx
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ---- palette (extracted from the Phase 00 deck) ---------------------------
GRAY   = RGBColor(0x60, 0x64, 0x6C)
BLACK  = RGBColor(0x00, 0x00, 0x00)
BLUE   = RGBColor(0x3D, 0x8D, 0xFF)
CARD   = RGBColor(0xF6, 0xF6, 0xF6)
BORDER = RGBColor(0xB8, 0xBC, 0xC4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
HILITE = RGBColor(0xD0, 0xED, 0xFA)

TITLE_FONT = "Calibri Light"
BODY_FONT  = "Calibri"
KICKER = "PHASE 02 QUALITY"

prs = Presentation()
prs.slide_width  = Emu(12192000)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]


# ---- primitives -----------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def txt(s, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
        font=BODY_FONT, anchor=None):
    tb = s.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.name = font; f.color.rgb = color
    return tb


def box(s, l, t, w, h, fill=CARD, border=BORDER, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Emu(l), Emu(t), Emu(w), Emu(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = border; sp.line.width = Emu(9525)
    sp.shadow.inherit = False
    return sp


def arrow(s, l, t, w, h, color=BLUE):
    sp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(l), Emu(t), Emu(w), Emu(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.color.rgb = color; sp.line.width = Emu(0)
    sp.shadow.inherit = False
    return sp


def hline(s, l, t, w, color=BORDER, weight=12700):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(l), Emu(t), Emu(l + w), Emu(t))
    cn.line.color.rgb = color; cn.line.width = Emu(weight)
    return cn


def kicker(s, text=KICKER):
    txt(s, 533400, 361950, 4000500, 247650, text, 11.25, GRAY, bold=True)


def title(s, text):
    txt(s, 533400, 876300, 10985400, 1219200, text, 37.5, BLACK, bold=True, font=TITLE_FONT)


def subhead(s, text):
    txt(s, 533400, 2209800, 9525000, 857250, text, 17.25, GRAY)


def pagenum(s, n):
    txt(s, 11372850, 6400800, 381000, 228600, f"{n:02d}", 9.75, GRAY, align=PP_ALIGN.RIGHT)


def emphasis(s, text):
    txt(s, 1676400, 5181600, 8839200, 600075, text, 21.0, BLACK, bold=True,
        align=PP_ALIGN.CENTER, font=TITLE_FONT, anchor=MSO_ANCHOR.MIDDLE)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def bullet_col(s, l, t, w, header, items, fill=CARD, header_size=22.5):
    h = 857250 + 333375 * len(items) + 285750
    box(s, l, t, w, h, fill=fill)
    txt(s, l + 266700, t + 250000, w - 533400, 381000, header, header_size, BLACK, bold=True)
    y = t + 857250
    for it in items:
        txt(s, l + 285750, y, 209550, 266700, "•", 15.75, BLUE, bold=True)
        txt(s, l + 609600, y, w - 876300, 323850, it, 15.75, BLACK)
        y += 333375
    return h


def chevron(s, cells, t=3314700, highlight_last=False):
    """cells: list of (header, sub) rendered as boxes joined by arrows."""
    n = len(cells)
    bw, aw, gap = 1981200, 342900, 133350
    total = n * bw + (n - 1) * (aw + 2 * gap)
    l = (12192000 - total) // 2
    h = 1504950
    for i, (header, sub) in enumerate(cells):
        fill = HILITE if (highlight_last and i == n - 1) else CARD
        box(s, l, t, bw, h, fill=fill)
        txt(s, l + 95250, t + 380000, bw - 190500, 400050, header, 20.0, BLACK,
            bold=True, align=PP_ALIGN.CENTER)
        if sub:
            txt(s, l + 95250, t + 800100, bw - 190500, 500000, sub, 13.5, GRAY,
                align=PP_ALIGN.CENTER)
        if i < n - 1:
            arrow(s, l + bw + gap, t + (h - 247650) // 2, aw, 247650, color=BORDER)
        l += bw + aw + 2 * gap


def pills(s, labels, t=3600450):
    n = len(labels)
    pw, aw, gap = 1676400, 200025, 76200
    total = n * pw + (n - 1) * (aw + 2 * gap)
    l = (12192000 - total) // 2
    h = 685800
    for i, lab in enumerate(labels):
        p = box(s, l, t, pw, h, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, l, t, pw, h, lab, 13.5, BLACK, bold=True, align=PP_ALIGN.CENTER,
            font=TITLE_FONT, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            arrow(s, l + pw + gap, t + (h - 200025) // 2, aw, 200025, color=BORDER)
        l += pw + aw + 2 * gap


# ==== SLIDE 1 — title ======================================================
s = slide()
txt(s, 533400, 1047750, 10985400, 762000, "Generation Got Cheap", 43.5, BLACK, bold=True, font=TITLE_FONT)
txt(s, 533400, 1828800, 10985400, 819150, "Review Is the Bottleneck", 43.5, BLACK, bold=True, font=TITLE_FONT)
txt(s, 533400, 3086100, 6000000, 361950, "Phase 02 quality engineering", 18.75, GRAY)
hline(s, 533400, 3943350, 7658100)
txt(s, 533400, 4343400, 8382000, 819150,
    "A green pipeline is necessary, not sufficient. Someone still has to understand what shipped.",
    21.75, BLACK)
pagenum(s, 1)
notes(s, "Land the turn: Phase 00 named the closing bracket; Phase 01 built something real. "
         "Today is not a new app — it is proving and understanding what already exists. "
         "Generation is cheap now; comprehension and review are the scarce skills. Frame before touching a terminal.")

# ==== SLIDE 2 — generation vs comprehension ================================
s = slide(); kicker(s)
title(s, "Phase 01 built it. Phase 02 proves you understand it.")
bullet_col(s, 533400, 3143250, 4972050, "Generation got cheap",
           ["Agents write more than you can read", "A green chat message", "4× the output"])
bullet_col(s, 6191250, 3143250, 4972050, "Comprehension stayed scarce",
           ["Reading and judgment", "A green CI run isn't understanding", "The review bottleneck"])
pagenum(s, 2)
notes(s, "The whole phase pivots here. Yesterday 'it compiles' felt like 'it's done'. "
         "Generation is abundant; comprehension is the leveraged, scarce skill. "
         "A green chat message is not proof; a green CI run is not understanding.")

# ==== SLIDE 3 — three kinds of broken ======================================
s = slide(); kicker(s)
title(s, "There are three kinds of broken")
cards3 = [
    ("Build / lint / test", "Machine-checkable. It won't compile or a test fails. Scanners find these."),
    ("Architecture", "Works, but wired wrong — auth in the wrong layer, secrets in client code. Invisible to tests."),
    ("Severity", "Not all findings are equal. A security gap outranks a lint warning."),
]
cw, gap = 3429000, 417900
lx = 533400
for header, body in cards3:
    box(s, lx, 3143250, cw, 1885950, fill=WHITE)
    txt(s, lx + 190500, 3390900, cw - 381000, 381000, header, 20.0, BLACK, bold=True)
    txt(s, lx + 190500, 3886200, cw - 381000, 1000000, body, 15.75, GRAY)
    lx += cw + gap
emphasis(s, "Scanners find only. You judge severity before anything gets fixed.")
pagenum(s, 3)
notes(s, "Three failure classes. The machine-checkable ones are easy; architecture violations are "
         "the dangerous ones — tests pass, wiring is wrong. Severity is judgment: why does a "
         "security gap outrank a lint warning?")

# ==== SLIDE 4 — scanner/fixer pairs ========================================
s = slide(); kicker(s)
title(s, "Scanners find. Fixers fix. Never the same step.")
pairs = [("/scan-errors", "/fix-errors"), ("/monitor", "/harden")]
row_t = 3143250
for scan, fix in pairs:
    box(s, 2667000, row_t, 2743200, 762000, fill=CARD)
    txt(s, 2667000, row_t, 2743200, 762000, scan, 18.0, BLACK, bold=True,
        align=PP_ALIGN.CENTER, font=TITLE_FONT, anchor=MSO_ANCHOR.MIDDLE)
    arrow(s, 5638800, row_t + 257175, 685800, 247650, color=BLUE)
    box(s, 6553200, row_t, 2743200, 762000, fill=CARD)
    txt(s, 6553200, row_t, 2743200, 762000, fix, 18.0, BLACK, bold=True,
        align=PP_ALIGN.CENTER, font=TITLE_FONT, anchor=MSO_ANCHOR.MIDDLE)
    row_t += 1000125
emphasis(s, "Scanners never change code. That separation is what makes the report trustworthy.")
pagenum(s, 4)
notes(s, "Two pipeline pairs. The scanner reports; you read and judge; then the fixer touches code. "
         "Auto-fix without a human reading the report is how you merge comprehension debt faster.")

# ==== SLIDE 5 — closing bracket is a discipline ============================
s = slide(); kicker(s)
title(s, "The closing bracket is a discipline, not a command")
chevron(s, [("Scan", "find"), ("Read", "judge"), ("Fix", "at the root"),
            ("Gate", "block if dirty"), ("Review", "compound")], highlight_last=True)
emphasis(s, "A green pipeline is necessary, not sufficient — someone still has to understand it.")
pagenum(s, 5)
notes(s, "The closing bracket is a chain, not one button: scan → read → fix → gate → "
         "review → compound. Interactive gates run when a human invokes them; the discipline is "
         "reading each report before the next step.")

# ==== SLIDE 6 — comprehension debt / bottleneck ============================
s = slide(); kicker(s)
title(s, "4× the output is not 4× the value")
bullet_col(s, 533400, 3143250, 4972050, "Old bottleneck",
           ["Can we build it?", "Typing code was the constraint", "Generation was expensive"])
bullet_col(s, 6191250, 3143250, 4972050, "New bottleneck",
           ["Can anyone understand it?", "Reading and judgment are the constraint", "Comprehension is expensive"])
emphasis(s, "Comprehension debt: when 'it works' and 'we understand it' drift apart.")
pagenum(s, 6)
notes(s, "4× output only becomes 4× value if someone reviewed it. Cite Osmani / DORA "
         "directionally — verify exact numbers before teaching. The gap between 'it works' and "
         "'we understand it' is the real risk.")

# ==== SLIDE 7 — ship like a professional ===================================
s = slide(); kicker(s)
title(s, "Ship like a professional")
quad = [
    ("Commit", "a checkpoint whose message explains why"),
    ("Quality gate", "build, tests, no violations — blocks until clean"),
    ("Pull request", "changes someone reviews before they join main"),
    ("Code review", "correct? safe? consistent? the durable human skill"),
]
positions = [(533400, 3143250), (6191250, 3143250), (533400, 4362450), (6191250, 4362450)]
for (header, body), (lx, ty) in zip(quad, positions):
    txt(s, lx, ty, 4972050, 381000, header, 21.0, BLACK, bold=True, font=TITLE_FONT)
    txt(s, lx, ty + 381000, 4972050, 685800, body, 16.5, GRAY)
pagenum(s, 7)
notes(s, "Define the shipping vocabulary precisely. The quality gate is the non-negotiable checklist "
         "that blocks a commit until it's clean. Code review — judging correctness, safety, and "
         "architecture fit — is the durable human skill, practiced as a process.")

# ==== SLIDE 8 — the loop that learns =======================================
s = slide(); kicker(s)
title(s, "The clever part: standards that compound")
chevron(s, [("Review finds a gap", ""), ("Rule lands in pr-requirements", ""),
            ("Every future commit checked", "")], t=3638550, highlight_last=True)
emphasis(s, "The bar rises by itself. The same mistake can't pass twice.")
pagenum(s, 8)
notes(s, "This is the compound-interest idea. When /pr-eval finds a new kind of problem, that rule is "
         "written into pr-requirements, so every future commit is checked for it. One line beats the "
         "same review comment every sprint. Explain why this matters more than any single fix.")

# ==== SLIDE 9 — committed rules vs personal memory =========================
s = slide(); kicker(s)
title(s, "Committed rules beat personal memory")
bullet_col(s, 533400, 3143250, 4972050, "Committed rules (git)",
           ["Survive a new hire's laptop", "Team truth, versioned", "Reviewed like code"])
bullet_col(s, 6191250, 3143250, 4972050, "Tool Memories (personal)",
           ["Local to one machine", "Not shared with the team", "Lost on reset"])
emphasis(s, "Team truth lives in git. Personal memory doesn't.")
pagenum(s, 9)
notes(s, "pr-requirements lives in the repo — it survives a new hire's laptop. Personal tool "
         "memories don't. Who approves org-wide rules is a process question, not a tooling one "
         "(adoption kit 04–05, docs/ai-program/).")

# ==== SLIDE 10 — observability =============================================
s = slide(); kicker(s)
title(s, "Observability lives in the diff, not the chat")
subhead(s, "The audit trail you'd show a customer is not a green chat message.")
items10 = [
    ("Diffs", "what actually changed"),
    ("Commits", "why it changed"),
    ("Test output", "proof it runs"),
    ("Pull requests", "who reviewed and approved"),
]
y = 3505200
for header, body in items10:
    txt(s, 819150, y, 209550, 266700, "•", 15.75, BLUE, bold=True)
    txt(s, 1143000, y, 3200000, 323850, header, 16.5, BLACK, bold=True)
    txt(s, 3600450, y, 6000000, 323850, body, 16.5, GRAY)
    y += 466725
pagenum(s, 10)
notes(s, "Ask: what audit trail would you show a customer — a chat log or the GitHub PR? "
         "Observability is the diff, the commit message, the test output, and the review — "
         "not a green message in the transcript.")

# ==== SLIDE 11 — nightly hygiene tiers =====================================
s = slide(); kicker(s)
title(s, "The closing bracket runs on a schedule too")
bullet_col(s, 533400, 3143250, 4972050, "Tier 1 — build · lint · test",
           ["Runs on a cron", "No agent API key", "Proves the repo still compiles"])
bullet_col(s, 6191250, 3143250, 4972050, "Tier 2 — scan → fix → monitor → harden",
           ["Opens a maintenance PR", "Humans still merge", "Needs boundaries + budget"], fill=HILITE)
emphasis(s, "Auto-merge to main is out of scope until your team documents the exception.")
pagenum(s, 11)
notes(s, "Interactive gates only run when invoked; nightly hygiene is the repo immune system. Tier 1 "
         "proves build/lint/test on a schedule with no key. Tier 2 opens a maintenance PR — humans "
         "still merge. Do not enable Tier 2 in the room without agreed decision boundaries and budget.")

# ==== SLIDE 12 — what this phase teaches ====================================
s = slide(); kicker(s)
title(s, "This phase teaches the closing bracket as a process")
quad2 = [
    ("Scan honestly", "classify before you fix"),
    ("Fix at the root", "not whack-a-mole symptoms"),
    ("Gate before ship", "clean, or it doesn't go"),
    ("Review to compound", "every gap becomes a rule"),
]
for (header, body), (lx, ty) in zip(quad2, positions):
    txt(s, lx, ty, 4972050, 381000, header, 21.0, BLACK, bold=True, font=TITLE_FONT)
    txt(s, lx, ty + 381000, 4972050, 685800, body, 16.5, GRAY)
emphasis(s, "We don't teach people to trust the pipeline. We teach them to prove what shipped.")
pagenum(s, 12)
notes(s, "Four moves become the operating model. The point is not blind trust in the pipeline — it's "
         "proving and understanding what shipped, and making standards compound so the next ship is "
         "harder to get wrong.")

# ==== SLIDE 13 — what runs in the lab ======================================
s = slide(); kicker(s)
title(s, "What runs in the lab")
pills(s, ["/scan-errors", "/fix-errors", "/monitor", "/harden", "/commit pr", "/pr-eval"])
emphasis(s, "Read each report before the fixer runs. Nobody watches the scroll.")
pagenum(s, 13)
notes(s, "The interactive chain. Kick it off, then teach while it runs — the scan is fast and the "
         "lesson is in the discussion, not the scroll. Read error-report.md and monitor-report.md "
         "before /fix-errors or /harden touches code.")

# ==== SLIDE 14 — deliverable ===============================================
s = slide(); kicker(s)
title(s, "The app gets cleaner. The discipline is the deliverable.")
bullet_col(s, 533400, 3143250, 4972050, "In the repo",
           ["Cleaner Phase 01 app", "Error + monitor reports", "Stronger pr-requirements", "Optional nightly Tier 1"])
bullet_col(s, 6191250, 3143250, 4972050, "In your head",
           ["The quality chain before every ship", "Gates are necessary, not sufficient",
            "Observability in diffs and PRs", "A path to org guardrails"])
pagenum(s, 14)
notes(s, "The cleaner app is the demonstration; the repeatable discipline is the deliverable. What they "
         "carry into every project: the quality chain, the belief that gates are necessary but not "
         "sufficient, and observability that lives in diffs, commits, and PRs.")

# ==== SLIDE 15 — transition / predict then run =============================
s = slide()
txt(s, 533400, 876300, 10985400, 1219200, "Now predict, then run.", 37.5, BLACK, bold=True, font=TITLE_FONT)
subhead(s, "Before the scanner runs, guess what will break in the app you were proud of yesterday.")
txt(s, 533400, 3143250, 4000500, 361950, "Prediction prompts", 18.75, GRAY, bold=True)
preds = [
    "What lint or type errors are hiding?",
    "What will /monitor find that the tests missed?",
    "Which AI-written function can nobody explain?",
    "What would you refuse to merge even if CI is green?",
]
y = 3695700
for p in preds:
    txt(s, 819150, y, 209550, 266700, "•", 15.75, BLUE, bold=True)
    txt(s, 1143000, y, 9500000, 361950, p, 17.25, BLACK)
    y += 466725
pagenum(s, 15)
notes(s, "Manufacture the discussion with predict-then-compare. Yesterday's 'done' becomes today's "
         "teaching moment. Write guesses on the board, then run the pipeline and compare against the "
         "reports. The reflection is the lesson.")

# ---- save -----------------------------------------------------------------
import os
out = "marketing/decks/phase-02-quality-engineering.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f"wrote {out} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
