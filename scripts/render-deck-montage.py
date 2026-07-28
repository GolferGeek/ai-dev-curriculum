#!/usr/bin/env python3
"""Rasterize a .pptx to per-slide PNGs + a montage, using PIL.

Reads the actual shapes/fills/text from the file (not the generator), so the
preview reflects what will open in PowerPoint and reveals text overflow.
Arial stands in for Calibri (close metrics) for on-screen preview only.

Usage: python3 scripts/render-deck-montage.py <deck.pptx> <out_dir> <montage.png>
"""
import sys, os
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

DECK = sys.argv[1] if len(sys.argv) > 1 else "marketing/decks/phase-02-quality-engineering.pptx"
OUT  = sys.argv[2] if len(sys.argv) > 2 else "outputs/phase-02-quality-engineering"
MONT = sys.argv[3] if len(sys.argv) > 3 else "outputs/phase-02-quality-engineering-montage.png"

prs = Presentation(DECK)
EMU_W, EMU_H = prs.slide_width, prs.slide_height
PXW, PXH = 1600, int(1600 * EMU_H / EMU_W)          # 1600x900
S = PXW / EMU_W                                       # EMU -> px
def px(v): return int(round(v * S))
def pt2px(pt): return int(round(pt / 72.0 * (PXW / (EMU_W / 914400.0))))  # dpi = PXW/inches

REG  = "/System/Library/Fonts/Supplemental/Arial.ttf"
BOLD = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"
_fc = {}
def font(pt, bold):
    key = (pt, bold)
    if key not in _fc:
        _fc[key] = ImageFont.truetype(BOLD if bold else REG, max(8, pt2px(pt)))
    return _fc[key]

NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
def prst_of(sh):
    try:
        g = sh._element.spPr.find(NS + "prstGeom")
        return g.get("prst") if g is not None else None
    except Exception:
        return None

def fill_rgb(sh):
    try:
        f = sh.fill
        if f.type is not None and f.fore_color.type is not None:
            return tuple(f.fore_color.rgb)
    except Exception:
        pass
    return None

def line_rgb(sh):
    try:
        ln = sh.line
        if ln.color and ln.color.type is not None:
            return tuple(ln.color.rgb)
    except Exception:
        pass
    return None

def run_props(sh):
    """First run's (size_pt, bold, color, align, anchor) for the shape."""
    tf = sh.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            sz = r.font.size.pt if r.font.size else 18
            bold = bool(r.font.bold)
            try:
                col = tuple(r.font.color.rgb) if r.font.color.type is not None else (0, 0, 0)
            except Exception:
                col = (0, 0, 0)
            align = str(p.alignment) if p.alignment is not None else "LEFT (1)"
            anc = str(tf.vertical_anchor) if tf.vertical_anchor is not None else "TOP"
            return sz, bold, col, align, anc
    return None

def wrap(draw, text, fnt, maxw):
    words, lines, cur = text.split(), [], ""
    for w in words:
        t = (cur + " " + w).strip()
        if draw.textlength(t, font=fnt) <= maxw or not cur:
            cur = t
        else:
            lines.append(cur); cur = w
    if cur:
        lines.append(cur)
    return lines

def draw_shape(d, sh):
    l, t, w, h = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
    prst = prst_of(sh)
    fill = fill_rgb(sh); line = line_rgb(sh)
    # backgrounds / connectors
    if prst == "rightArrow":
        c = fill or (61, 141, 255)
        shaft = h * 0.42
        tip = l + int(w * 0.55)
        y0, y1 = t + (h - shaft) / 2, t + (h + shaft) / 2
        d.polygon([(l, y0), (tip, y0), (tip, t), (l + w, t + h / 2),
                   (tip, t + h), (tip, y1), (l, y1)], fill=c)
        return
    if prst in ("rect", "roundRect") and fill is not None:
        r = 14 if prst == "roundRect" else 0
        d.rounded_rectangle([l, t, l + w, t + h], radius=r, fill=fill,
                            outline=line, width=2 if line else 0)
    elif sh.shape_type is not None and str(sh.shape_type).startswith("LINE") or prst == "line":
        c = line or (184, 188, 196)
        if h == 0:
            d.line([(l, t), (l + w, t)], fill=c, width=3)
        else:
            d.line([(l, t), (l, t + h)], fill=c, width=3)
    # connectors added via add_connector show as shape_type None but have a line
    elif sh.width == 0 or sh.height == 0:
        c = line or (184, 188, 196)
        d.line([(l, t), (l + w, t + h)], fill=c, width=3)

def draw_text(d, sh):
    if not sh.has_text_frame or not sh.text_frame.text.strip():
        return
    rp = run_props(sh)
    if not rp:
        return
    sz, bold, col, align, anc = rp
    fnt = font(sz, bold)
    l, t, w, h = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
    text = sh.text_frame.text
    lines = []
    for para in text.split("\n"):
        lines += wrap(d, para, fnt, w) if para else [""]
    lh = int(fnt.size * 1.2)
    total = lh * len(lines)
    y = t + (h - total) // 2 if "MIDDLE" in anc else t
    for ln in lines:
        tw = d.textlength(ln, font=fnt)
        if "CENTER" in align:
            x = l + (w - tw) // 2
        elif "RIGHT" in align:
            x = l + w - tw
        else:
            x = l
        d.text((x, y), ln, font=fnt, fill=col)
        y += lh

os.makedirs(OUT, exist_ok=True)
imgs = []
for i, s in enumerate(prs.slides, 1):
    img = Image.new("RGB", (PXW, PXH), "white")
    d = ImageDraw.Draw(img)
    for sh in s.shapes:          # backgrounds first
        draw_shape(d, sh)
    for sh in s.shapes:          # text on top
        draw_text(d, sh)
    d.rectangle([0, 0, PXW - 1, PXH - 1], outline=(220, 222, 226), width=1)
    p = os.path.join(OUT, f"slide-{i}.png")
    img.save(p)
    imgs.append(img)

# montage: 5 cols x 3 rows
cols, rows, pad = 5, 3, 16
tw, th = 400, int(400 * PXH / PXW)
mont = Image.new("RGB", (cols * tw + (cols + 1) * pad, rows * th + (rows + 1) * pad), (245, 246, 248))
for idx, im in enumerate(imgs):
    r, c = divmod(idx, cols)
    x = pad + c * (tw + pad)
    y = pad + r * (th + pad)
    mont.paste(im.resize((tw, th)), (x, y))
os.makedirs(os.path.dirname(MONT) or ".", exist_ok=True)
mont.save(MONT)
print(f"wrote {len(imgs)} slide PNGs to {OUT}/ and montage {MONT} ({mont.size[0]}x{mont.size[1]})")
